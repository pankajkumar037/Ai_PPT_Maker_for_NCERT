"""
HTML to PowerPoint Converter
Converts HTML presentations to PowerPoint with matching styling
"""

import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import re

# Import Pexels image fetcher
try:
    from images.pexels_fetcher import get_image_for_slide
    PEXELS_AVAILABLE = True
except ImportError:
    PEXELS_AVAILABLE = False
    print("Warning: Pexels image fetcher not available")


def convert_html_style_to_pptx(slides_data, output_filename, html_style="vibrant"):
    """
    Creates a PowerPoint presentation matching HTML template styling

    Args:
        slides_data: List of slide dictionaries
        output_filename: Name for output PPTX file
        html_style: "vibrant", "modern", or "dark"

    Returns:
        Path to the created PowerPoint file
    """

    # Create new presentation
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # Define color schemes matching HTML templates
    if html_style == "vibrant":
        # Purple to Blue gradient
        title_gradient = [
            RGBColor(109, 40, 217),   # Purple-600
            RGBColor(59, 130, 246),    # Blue-500
            RGBColor(34, 211, 238)     # Cyan-400
        ]
        header_color1 = RGBColor(109, 40, 217)  # Purple-600
        header_color2 = RGBColor(59, 130, 246)   # Blue-500
        circle_color = RGBColor(249, 115, 22)    # Orange-500
        accent_color = RGBColor(249, 115, 22)    # Orange-500
        text_color = RGBColor(31, 41, 55)        # Gray-800
        bg_gradient = [
            RGBColor(249, 250, 251),   # Gray-50
            RGBColor(243, 244, 246)    # Gray-100
        ]
    elif html_style == "modern":
        # Blue to Indigo gradient
        title_gradient = [
            RGBColor(37, 99, 235),     # Blue-600
            RGBColor(59, 130, 246),    # Blue-500
            RGBColor(96, 165, 250)     # Blue-400
        ]
        header_color1 = RGBColor(37, 99, 235)    # Blue-600
        header_color2 = RGBColor(59, 130, 246)   # Blue-500
        circle_color = RGBColor(59, 130, 246)    # Blue-500
        accent_color = RGBColor(59, 130, 246)    # Blue-600
        text_color = RGBColor(17, 24, 39)        # Gray-900
        bg_gradient = [
            RGBColor(249, 250, 251),   # Gray-50
            RGBColor(243, 244, 246)    # Gray-100
        ]
    elif html_style == "dark":
        # Dark gradient
        title_gradient = [
            RGBColor(17, 24, 39),      # Gray-900
            RGBColor(31, 41, 55),      # Gray-800
            RGBColor(55, 65, 81)       # Gray-700
        ]
        header_color1 = RGBColor(17, 24, 39)     # Gray-900
        header_color2 = RGBColor(31, 41, 55)     # Gray-800
        circle_color = RGBColor(251, 191, 36)    # Yellow-400
        accent_color = RGBColor(251, 191, 36)    # Yellow-400
        text_color = RGBColor(229, 231, 235)     # Gray-200
        bg_gradient = [
            RGBColor(31, 41, 55),      # Gray-800
            RGBColor(17, 24, 39)       # Gray-900
        ]
    else:
        # Default to vibrant
        title_gradient = [RGBColor(109, 40, 217), RGBColor(59, 130, 246), RGBColor(34, 211, 238)]
        header_color1 = RGBColor(109, 40, 217)
        header_color2 = RGBColor(59, 130, 246)
        circle_color = RGBColor(249, 115, 22)
        accent_color = RGBColor(249, 115, 22)
        text_color = RGBColor(31, 41, 55)
        bg_gradient = [RGBColor(249, 250, 251), RGBColor(243, 244, 246)]

    # Create slides
    for i, slide_data in enumerate(slides_data):
        is_title = i == 0

        # Add blank slide
        slide = prs.slides.add_slide(prs.slide_layouts[6])

        if is_title:
            # ===== TITLE SLIDE =====
            # Gradient background (3 layers)
            for idx, color in enumerate(title_gradient):
                rect = slide.shapes.add_shape(
                    1,  # Rectangle
                    Inches(0), Inches(idx * 2.5), Inches(10), Inches(2.5)
                )
                rect.fill.solid()
                rect.fill.fore_color.rgb = color
                rect.line.fill.background()

            # Decorative circles
            circles = [
                (8.5, 0.5, 1.5, circle_color),
                (0.8, 5.5, 1.2, circle_color),
                (9, 6, 1.0, circle_color)
            ]
            for x, y, size, color in circles:
                circle = slide.shapes.add_shape(9, Inches(x), Inches(y), Inches(size), Inches(size))
                circle.fill.solid()
                circle.fill.fore_color.rgb = color
                circle.line.fill.background()

            # Badge
            badge = slide.shapes.add_shape(1, Inches(3.5), Inches(1.8), Inches(3), Inches(0.4))
            badge.fill.solid()
            badge.fill.fore_color.rgb = RGBColor(255, 255, 255)
            badge.fill.transparency = 0.7
            badge.line.fill.background()

            badge_frame = badge.text_frame
            badge_frame.text = "✨ AI-Powered Presentation"
            badge_para = badge_frame.paragraphs[0]
            badge_para.font.size = Pt(12)
            badge_para.font.bold = True
            badge_para.font.color.rgb = RGBColor(255, 255, 255) if html_style != "dark" else RGBColor(17, 24, 39)
            badge_para.alignment = PP_ALIGN.CENTER

            # Title
            title_box = slide.shapes.add_textbox(Inches(1), Inches(2.8), Inches(8), Inches(2))
            title_frame = title_box.text_frame
            title_frame.text = slide_data['title']
            title_frame.word_wrap = True
            title_para = title_frame.paragraphs[0]
            title_para.font.size = Pt(54)
            title_para.font.bold = True
            title_para.font.color.rgb = RGBColor(255, 255, 255)
            title_para.alignment = PP_ALIGN.CENTER
            title_para.line_spacing = 1.2

            # Decorative line
            line = slide.shapes.add_shape(1, Inches(4.2), Inches(5.0), Inches(1.6), Inches(0.06))
            line.fill.solid()
            line.fill.fore_color.rgb = RGBColor(255, 255, 255)
            line.fill.transparency = 0.4
            line.line.fill.background()

            # Subtitle
            if slide_data.get('content') and len(slide_data['content']) > 0:
                subtitle_box = slide.shapes.add_textbox(Inches(1), Inches(5.3), Inches(8), Inches(1.2))
                subtitle_frame = subtitle_box.text_frame
                subtitle_frame.text = slide_data['content'][0]
                subtitle_frame.word_wrap = True
                subtitle_para = subtitle_frame.paragraphs[0]
                subtitle_para.font.size = Pt(22)
                subtitle_para.font.color.rgb = RGBColor(255, 255, 255) if html_style != "dark" else RGBColor(229, 231, 235)
                subtitle_para.alignment = PP_ALIGN.CENTER
                subtitle_para.line_spacing = 1.3

        else:
            # ===== CONTENT SLIDE =====
            # Background gradient
            for idx, color in enumerate(bg_gradient):
                rect = slide.shapes.add_shape(
                    1, Inches(0), Inches(idx * 3.75), Inches(10), Inches(3.75)
                )
                rect.fill.solid()
                rect.fill.fore_color.rgb = color
                rect.line.fill.background()

            # Header gradient
            header_parts = [
                (header_color1, 0.4, 0.4, 4.6, 1.1),
                (header_color2, 5, 0.4, 4.6, 1.1)
            ]
            for color, x, y, w, h in header_parts:
                rect = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(h))
                rect.fill.solid()
                rect.fill.fore_color.rgb = color
                rect.line.fill.background()

            # Accent bar
            accent_bar = slide.shapes.add_shape(1, Inches(0.4), Inches(1.6), Inches(1.2), Inches(0.08))
            accent_bar.fill.solid()
            accent_bar.fill.fore_color.rgb = accent_color
            accent_bar.line.fill.background()

            # Content card
            card_color = RGBColor(55, 65, 81) if html_style == "dark" else RGBColor(255, 255, 255)
            content_card = slide.shapes.add_shape(1, Inches(0.4), Inches(1.9), Inches(9.2), Inches(5.4))
            content_card.fill.solid()
            content_card.fill.fore_color.rgb = card_color
            card_border = RGBColor(75, 85, 99) if html_style == "dark" else RGBColor(209, 213, 219)
            content_card.line.color.rgb = card_border
            content_card.line.width = Pt(1)

            # Title in header
            title_box = slide.shapes.add_textbox(Inches(0.6), Inches(0.6), Inches(8.8), Inches(0.7))
            title_frame = title_box.text_frame
            title_frame.text = slide_data['title']
            title_para = title_frame.paragraphs[0]
            title_para.font.size = Pt(36)
            title_para.font.bold = True
            title_para.font.color.rgb = RGBColor(255, 255, 255)

            # Content
            has_image = slide_data.get('has_image', False)

            if has_image:
                content_left = Inches(0.7)
                content_width = Inches(4.3)
                image_left = Inches(5.3)
                image_width = Inches(3.9)
                image_top = Inches(2.1)
                image_height = Inches(4.9)
            else:
                content_left = Inches(0.7)
                content_width = Inches(8.6)

            # Bullet points
            content_top = Inches(2.1)
            content_height = Inches(4.9)

            content_box = slide.shapes.add_textbox(content_left, content_top, content_width, content_height)
            text_frame = content_box.text_frame
            text_frame.word_wrap = True
            text_frame.auto_size = None  # Disable auto-size to prevent overflow
            text_frame.margin_left = Inches(0.15)
            text_frame.margin_right = Inches(0.15)
            text_frame.margin_top = Inches(0.1)
            text_frame.margin_bottom = Inches(0.1)

            for j, point in enumerate(slide_data.get('content', [])):
                if j == 0:
                    para = text_frame.paragraphs[0]
                else:
                    para = text_frame.add_paragraph()

                para.level = 0
                para.space_before = Pt(8)
                para.space_after = Pt(6)
                para.line_spacing = 1.15  # Tighter line height to prevent overflow

                # Parse bold formatting
                parts = re.split(r'\*\*(.*?)\*\*', point)

                for idx, part in enumerate(parts):
                    if not part:
                        continue
                    run = para.add_run()
                    run.text = part
                    run.font.name = 'Segoe UI'
                    run.font.size = Pt(14)  # Reduced from 16
                    run.font.color.rgb = text_color

                    if idx % 2 == 1:  # Bold text
                        run.font.bold = True
                        run.font.size = Pt(16)  # Reduced from 18
                        if html_style != "dark":
                            run.font.color.rgb = accent_color

            # Add image (from Pexels or placeholder)
            if has_image:
                image_added = False

                # Try to fetch and add real image from Pexels
                if PEXELS_AVAILABLE:
                    try:
                        # Get image for this slide
                        image_path = get_image_for_slide(slide_data['title'], i)

                        if image_path and os.path.exists(image_path):
                            # Add the actual image
                            slide.shapes.add_picture(
                                image_path,
                                image_left, image_top,
                                width=image_width, height=image_height
                            )
                            image_added = True
                            print(f"  ✓ Added image from Pexels for slide {i+1}")
                    except Exception as e:
                        print(f"  ⚠ Could not add Pexels image: {e}")

                # Fallback to placeholder if image wasn't added
                if not image_added:
                    placeholder_gradient = [
                        (RGBColor(224, 231, 255), image_left, image_top, image_width, image_height/3),
                        (RGBColor(221, 214, 254), image_left, image_top + image_height/3, image_width, image_height/3),
                        (RGBColor(251, 207, 232), image_left, image_top + 2*image_height/3, image_width, image_height/3)
                    ]
                    for color, x, y, w, h in placeholder_gradient:
                        rect = slide.shapes.add_shape(1, x, y, w, h)
                        rect.fill.solid()
                        rect.fill.fore_color.rgb = color
                        rect.line.fill.background()

                    border = slide.shapes.add_shape(1, image_left, image_top, image_width, image_height)
                    border.fill.background()
                    border.line.color.rgb = RGBColor(99, 102, 241)
                    border.line.width = Pt(3)
                    border.line.dash_style = 2

                    icon_box = slide.shapes.add_textbox(
                        image_left + Inches(0.5),
                        image_top + image_height/2 - Inches(0.4),
                        image_width - Inches(1),
                        Inches(0.8)
                    )
                    icon_frame = icon_box.text_frame
                    icon_frame.text = "📸 Visual Content\nAdd image or diagram"
                    icon_para = icon_frame.paragraphs[0]
                    icon_para.font.size = Pt(16)
                    icon_para.font.bold = True
                    icon_para.font.color.rgb = RGBColor(79, 70, 229)
                    icon_para.alignment = PP_ALIGN.CENTER
                    icon_para.line_spacing = 1.4

        # Add speaker notes
        notes_slide = slide.notes_slide
        notes_frame = notes_slide.notes_text_frame
        notes_frame.text = slide_data.get('notes', '')

    # Save
    os.makedirs('output', exist_ok=True)
    output_path = f"output/{output_filename}"
    prs.save(output_path)

    print(f"PowerPoint with {html_style} styling saved: {output_path}")
    return output_path
