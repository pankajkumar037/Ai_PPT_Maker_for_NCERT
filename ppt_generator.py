"""
AI PPT Generator - Main Logic
This file contains all the functions needed to generate PowerPoint presentations using OpenAI.

Flow:
1. User provides a topic/prompt
2. OpenAI generates structured slide content
3. We parse the content and create PPT slides
4. Download relevant images from Pexels and add them to slides
"""

import os
from openai import OpenAI
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from dotenv import load_dotenv
import json

# ========== STEP 1: LOAD ENVIRONMENT VARIABLES ==========
# This loads your OpenAI API key from the .env file
load_dotenv()

# Import Pexels image fetcher
try:
    from images.pexels_fetcher import get_image_for_slide
    PEXELS_AVAILABLE = True
except ImportError:
    PEXELS_AVAILABLE = False
    print("Warning: Pexels image fetcher not available")


def initialize_openai_client():
    """
    Creates and returns an OpenAI client using the API key from .env file

    Returns:
        OpenAI client object ready to make API calls
    """
    # Get the API key from environment variables
    api_key = os.getenv('OPENAI_API_KEY')

    if not api_key:
        raise ValueError("OpenAI API key not found! Please add it to your .env file")

    # Set the API key in environment and create client - simple approach
    os.environ['OPENAI_API_KEY'] = api_key
    client = OpenAI()

    return client


def generate_slide_content(client, topic, num_slides=8):
    """
    Uses OpenAI to generate structured content for PowerPoint slides.

    Args:
        client: OpenAI client object
        topic: The topic/prompt provided by the user
        num_slides: Number of slides to generate (default: 8)

    Returns:
        List of dictionaries, each containing slide information:
        - title: Slide title
        - content: List of bullet points
        - notes: Speaker notes
        - has_image: Whether this slide should have an image placeholder
    """

    # Create a detailed prompt for OpenAI to generate structured content
    prompt = f"""
Create a COMPREHENSIVE and DETAILED educational PowerPoint presentation on the topic: "{topic}"

This presentation is for TEACHERS to use in classroom teaching for students (Class 1-12).
Make it THOROUGH, INFORMATIVE, and EASY TO UNDERSTAND.

Create {num_slides} slides with the following structure:

SLIDE 1 (Title Slide):
- Main title: The topic name
- Subtitle: A brief catchy description
- has_image: true

REMAINING SLIDES (Content Slides):
For each content slide, provide:
1. An ENGAGING and DESCRIPTIVE title
2. 4-6 WELL-FORMATTED bullet points
   - Keep each point concise but informative (2-3 lines max)
   - Include key definitions, examples, or explanations
   - Use simple language suitable for students
   - Format important terms by surrounding with ** for bold (e.g., "**Photosynthesis** is the process...")
3. Comprehensive speaker notes (4-6 sentences) with additional teaching tips
4. has_image: true for at least 50% of slides

IMPORTANT GUIDELINES:
- Cover ALL key aspects of the topic across multiple slides
- Include definitions, processes, examples, real-world applications
- Make content detailed enough for a 30-40 minute lecture
- Keep bullet points readable - not too much text on one slide
- Use **bold** for key terms and important concepts

Format your response as a JSON array where each slide is an object with these keys:
- "title": string (engaging and descriptive)
- "content": array of strings (4-6 formatted bullet points)
- "notes": string (comprehensive speaker notes with teaching tips)
- "has_image": boolean

Example of good formatted bullet point:
❌ BAD: "Plants need water for photosynthesis - it acts as a raw material and also helps transport nutrients from roots to leaves through the xylem vessels"
✅ GOOD: "**Water as Raw Material**: Essential for photosynthesis and nutrient transport through xylem vessels"
"""

    # Call OpenAI API to generate the content
    print("Calling OpenAI API to generate slide content...")
    response = client.chat.completions.create(
        model="gpt-4o-2024-08-06",
        messages=[
            {"role": "system", "content": "You are an expert educational content creator and experienced teacher (10+ years) who creates comprehensive, detailed, and engaging presentation content for classroom teaching (Class 1-12). Your presentations are thorough, well-structured, and perfect for 30-45 minute lectures."},
            {"role": "user", "content": prompt}
        ],
        temperature=0.4,
        max_tokens=8000
    )

    # Extract the generated text from the API response
    generated_text = response.choices[0].message.content

    # Parse the JSON response from OpenAI
    # Sometimes the response has markdown code blocks, so we clean it first
    if "```json" in generated_text:
        # Extract JSON from markdown code block
        generated_text = generated_text.split("```json")[1].split("```")[0].strip()
    elif "```" in generated_text:
        generated_text = generated_text.split("```")[1].split("```")[0].strip()

    try:
        slides_data = json.loads(generated_text)
        print(f"Successfully generated content for {len(slides_data)} slides!")
        return slides_data
    except json.JSONDecodeError as e:
        print(f"Error parsing OpenAI response: {e}")
        print("Raw response:", generated_text)
        raise


def create_presentation_from_template(template_name, slides_data, output_filename):
    """
    Creates a PowerPoint presentation using a template and the generated content.

    Args:
        template_name: Name of the template to use ("modern_blue" or "elegant_dark")
        slides_data: List of slide data dictionaries from generate_slide_content()
        output_filename: Name for the output PPT file (e.g., "my_presentation.pptx")

    Returns:
        Path to the created presentation file
    """

    # Load the selected template
    template_path = f"templates/{template_name}.pptx"

    if not os.path.exists(template_path):
        raise FileNotFoundError(f"Template not found: {template_path}")

    print(f"Loading template: {template_name}")
    prs = Presentation(template_path)

    # Define color schemes for each template
    if template_name == "ocean_breeze":
        text_color = RGBColor(30, 41, 59)       # Slate-800
        accent_color = RGBColor(251, 191, 36)   # Amber-400
        bg_color = RGBColor(255, 255, 255)      # White
    elif template_name == "sunset_glow":
        text_color = RGBColor(51, 65, 85)       # Slate-700
        accent_color = RGBColor(253, 224, 71)   # Yellow-300
        bg_color = RGBColor(255, 255, 255)      # White
    elif template_name == "forest_fresh":
        text_color = RGBColor(6, 95, 70)        # Emerald-800
        accent_color = RGBColor(252, 211, 77)   # Amber-300
        bg_color = RGBColor(255, 255, 255)      # White
    elif template_name == "royal_purple":
        text_color = RGBColor(55, 48, 163)      # Violet-800
        accent_color = RGBColor(251, 191, 36)   # Amber-400
        bg_color = RGBColor(255, 255, 255)      # White
    elif template_name == "corporate_slate":
        text_color = RGBColor(30, 41, 59)       # Slate-800
        accent_color = RGBColor(34, 211, 238)   # Cyan-400
        bg_color = RGBColor(255, 255, 255)      # White
    elif template_name == "modern_blue":
        text_color = RGBColor(26, 35, 126)      # Dark blue
        accent_color = RGBColor(33, 150, 243)   # Bright blue
        bg_color = RGBColor(255, 255, 255)      # White
    elif template_name == "elegant_dark":
        text_color = RGBColor(250, 250, 250)    # Off-white
        accent_color = RGBColor(255, 193, 7)    # Gold
        bg_color = RGBColor(18, 18, 18)         # Almost black
    elif template_name == "vibrant_education":
        text_color = RGBColor(44, 62, 80)       # Dark text
        accent_color = RGBColor(255, 140, 0)    # Bright orange
        bg_color = RGBColor(255, 255, 255)      # White
    else:
        # Default colors if template name doesn't match
        text_color = RGBColor(0, 0, 0)
        accent_color = RGBColor(0, 120, 215)
        bg_color = RGBColor(255, 255, 255)

    # Clear existing slides from template (we'll create new ones)
    # Keep the layouts but remove the sample slides
    while len(prs.slides) > 0:
        rId = prs.slides._sldIdLst[0].rId
        prs.part.drop_rel(rId)
        del prs.slides._sldIdLst[0]

    # Create slides based on the generated content
    for i, slide_data in enumerate(slides_data):
        print(f"Creating slide {i+1}: {slide_data['title']}")

        # Add a blank slide
        slide = prs.slides.add_slide(prs.slide_layouts[6])  # 6 = blank layout

        # Determine if this is the title slide (first slide)
        is_title_slide = (i == 0)

        if is_title_slide:
            # ===== ENHANCED TITLE SLIDE LAYOUT =====
            # Modern, professional designs matching HTML styling
            if template_name == "modern_blue":
                # Full gradient background
                bg_gradient = [
                    (RGBColor(37, 99, 235), 0, 0, 10, 2.5),    # Blue-600
                    (RGBColor(59, 130, 246), 0, 2.5, 10, 2.5),  # Blue-500
                    (RGBColor(96, 165, 250), 0, 5, 10, 2.5)     # Blue-400
                ]
                for color, x, y, w, h in bg_gradient:
                    rect = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(h))
                    rect.fill.solid()
                    rect.fill.fore_color.rgb = color
                    rect.line.fill.background()

                # Decorative circles
                circles = [
                    (8.5, 0.5, 1.5, RGBColor(147, 197, 253)),  # Blue-300
                    (0.8, 5.5, 1.2, RGBColor(191, 219, 254)),  # Blue-200
                    (9, 6, 1.0, RGBColor(59, 130, 246))        # Blue-500
                ]
                for x, y, size, color in circles:
                    circle = slide.shapes.add_shape(9, Inches(x), Inches(y), Inches(size), Inches(size))
                    circle.fill.solid()
                    circle.fill.fore_color.rgb = color
                    circle.line.fill.background()

                title_color = RGBColor(255, 255, 255)
                subtitle_color = RGBColor(255, 255, 255)

            elif template_name == "elegant_dark":
                # Dark gradient background
                bg_gradient = [
                    (RGBColor(17, 24, 39), 0, 0, 10, 2.5),    # Gray-900
                    (RGBColor(31, 41, 55), 0, 2.5, 10, 2.5),  # Gray-800
                    (RGBColor(55, 65, 81), 0, 5, 10, 2.5)     # Gray-700
                ]
                for color, x, y, w, h in bg_gradient:
                    rect = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(h))
                    rect.fill.solid()
                    rect.fill.fore_color.rgb = color
                    rect.line.fill.background()

                # Yellow/Gold accent circles
                circles = [
                    (8.5, 0.5, 1.5, RGBColor(251, 191, 36)),  # Yellow-400
                    (0.8, 5.5, 1.2, RGBColor(252, 211, 77)),  # Yellow-300
                    (9.2, 6.5, 0.8, RGBColor(250, 204, 21))   # Yellow-400
                ]
                for x, y, size, color in circles:
                    circle = slide.shapes.add_shape(9, Inches(x), Inches(y), Inches(size), Inches(size))
                    circle.fill.solid()
                    circle.fill.fore_color.rgb = color
                    circle.line.fill.background()

                title_color = RGBColor(255, 255, 255)
                subtitle_color = RGBColor(229, 231, 235)  # Gray-200

            elif template_name == "vibrant_education":
                # Colorful gradient background
                colors = [RGBColor(138, 43, 226), RGBColor(52, 152, 219), RGBColor(46, 204, 113)]
                for idx, color in enumerate(colors):
                    rect = slide.shapes.add_shape(1, Inches(0), Inches(idx * 2.5), Inches(10), Inches(2.5))
                    rect.fill.solid()
                    rect.fill.fore_color.rgb = color
                    rect.line.fill.background()

                # Decorative circles
                circle_data = [(0.5, 0.5, RGBColor(255, 215, 0)), (9, 1, RGBColor(255, 140, 0)),
                               (1, 6.5, RGBColor(46, 204, 113)), (8.5, 6, RGBColor(255, 215, 0))]
                for x, y, color in circle_data:
                    circle = slide.shapes.add_shape(9, Inches(x), Inches(y), Inches(0.8), Inches(0.8))
                    circle.fill.solid()
                    circle.fill.fore_color.rgb = color
                    circle.line.fill.background()

                title_color = RGBColor(255, 255, 255)
                subtitle_color = RGBColor(255, 255, 255)

            elif template_name in ["ocean_breeze", "sunset_glow", "forest_fresh", "royal_purple"]:
                # Generic gradient background for new templates
                if template_name == "ocean_breeze":
                    gradient_colors = [RGBColor(2, 132, 199), RGBColor(6, 182, 212), RGBColor(34, 211, 238)]
                    circle_color = RGBColor(251, 191, 36)
                elif template_name == "sunset_glow":
                    gradient_colors = [RGBColor(249, 115, 22), RGBColor(251, 146, 60), RGBColor(236, 72, 153)]
                    circle_color = RGBColor(253, 224, 71)
                elif template_name == "forest_fresh":
                    gradient_colors = [RGBColor(6, 95, 70), RGBColor(5, 150, 105), RGBColor(16, 185, 129)]
                    circle_color = RGBColor(252, 211, 77)
                elif template_name == "royal_purple":
                    gradient_colors = [RGBColor(109, 40, 217), RGBColor(139, 92, 246), RGBColor(167, 139, 250)]
                    circle_color = RGBColor(251, 191, 36)

                for idx, color in enumerate(gradient_colors):
                    rect = slide.shapes.add_shape(1, Inches(0), Inches(idx * 2.5), Inches(10), Inches(2.5))
                    rect.fill.solid()
                    rect.fill.fore_color.rgb = color
                    rect.line.fill.background()

                # Decorative circles
                circles = [(8.5, 0.5, 1.2), (0.8, 5.5, 0.8), (9, 6, 0.6)]
                for x, y, size in circles:
                    circle = slide.shapes.add_shape(9, Inches(x), Inches(y), Inches(size), Inches(size))
                    circle.fill.solid()
                    circle.fill.fore_color.rgb = circle_color
                    circle.line.fill.background()

                title_color = RGBColor(255, 255, 255)
                subtitle_color = RGBColor(255, 255, 255)

            elif template_name == "corporate_slate":
                # Split design
                left = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(6), Inches(7.5))
                left.fill.solid()
                left.fill.fore_color.rgb = RGBColor(30, 41, 59)
                left.line.fill.background()

                right = slide.shapes.add_shape(1, Inches(6), Inches(0), Inches(4), Inches(7.5))
                right.fill.solid()
                right.fill.fore_color.rgb = RGBColor(59, 130, 246)
                right.line.fill.background()

                # Accent line
                accent_line = slide.shapes.add_shape(1, Inches(5.8), Inches(0), Inches(0.4), Inches(7.5))
                accent_line.fill.solid()
                accent_line.fill.fore_color.rgb = accent_color
                accent_line.line.fill.background()

                title_color = RGBColor(255, 255, 255)
                subtitle_color = RGBColor(255, 255, 255)

            else:
                # Default fallback for unknown templates
                bg = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(10), Inches(7.5))
                bg.fill.solid()
                bg.fill.fore_color.rgb = RGBColor(59, 130, 246)
                bg.line.fill.background()

                title_color = RGBColor(255, 255, 255)
                subtitle_color = RGBColor(255, 255, 255)

            # Add modern badge (AI-Powered Presentation)
            badge = slide.shapes.add_shape(
                1,  # Rectangle
                Inches(3.5), Inches(1.8), Inches(3), Inches(0.4)
            )
            badge.fill.solid()
            badge.fill.fore_color.rgb = RGBColor(255, 255, 255)
            badge.fill.transparency = 0.7  # Semi-transparent
            badge.line.fill.background()

            badge_frame = badge.text_frame
            badge_frame.text = "✨ AI-Powered Presentation"
            badge_para = badge_frame.paragraphs[0]
            badge_para.font.size = Pt(12)
            badge_para.font.bold = True
            badge_para.font.color.rgb = RGBColor(255, 255, 255) if template_name != "elegant_dark" else RGBColor(17, 24, 39)
            badge_para.alignment = PP_ALIGN.CENTER

            # Add title text box - larger and centered
            title_box = slide.shapes.add_textbox(Inches(1), Inches(2.8), Inches(8), Inches(2))
            title_frame = title_box.text_frame
            title_frame.text = slide_data['title']
            title_frame.word_wrap = True
            title_para = title_frame.paragraphs[0]
            title_para.font.size = Pt(54)
            title_para.font.bold = True
            title_para.font.color.rgb = title_color
            title_para.alignment = PP_ALIGN.CENTER
            title_para.line_spacing = 1.2

            # Add decorative line
            line = slide.shapes.add_shape(1, Inches(4.2), Inches(5.0), Inches(1.6), Inches(0.06))
            line.fill.solid()
            line.fill.fore_color.rgb = RGBColor(255, 255, 255)
            line.fill.transparency = 0.4
            line.line.fill.background()

            # Add subtitle if present in content
            if slide_data.get('content') and len(slide_data['content']) > 0:
                subtitle_box = slide.shapes.add_textbox(Inches(1), Inches(5.3), Inches(8), Inches(1.2))
                subtitle_frame = subtitle_box.text_frame
                subtitle_frame.text = slide_data['content'][0]
                subtitle_frame.word_wrap = True
                subtitle_para = subtitle_frame.paragraphs[0]
                subtitle_para.font.size = Pt(22)
                subtitle_para.font.color.rgb = subtitle_color
                subtitle_para.alignment = PP_ALIGN.CENTER
                subtitle_para.line_spacing = 1.3

        else:
            # ===== ENHANCED CONTENT SLIDE LAYOUT =====
            # Modern card-style designs matching HTML
            if template_name == "modern_blue":
                # Light gradient background
                bg_gradient = [
                    (RGBColor(249, 250, 251), 0, 0, 10, 3.75),    # Gray-50
                    (RGBColor(243, 244, 246), 0, 3.75, 10, 3.75)  # Gray-100
                ]
                for color, x, y, w, h in bg_gradient:
                    rect = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(h))
                    rect.fill.solid()
                    rect.fill.fore_color.rgb = color
                    rect.line.fill.background()

                # Modern header card with gradient
                header_gradient = [
                    (RGBColor(37, 99, 235), 0.4, 0.4, 4.6, 1.1),   # Blue-600
                    (RGBColor(59, 130, 246), 5, 0.4, 4.6, 1.1)     # Blue-500
                ]
                for color, x, y, w, h in header_gradient:
                    rect = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(h))
                    rect.fill.solid()
                    rect.fill.fore_color.rgb = color
                    rect.line.fill.background()

                # Decorative accent bar under header
                accent_bar = slide.shapes.add_shape(1, Inches(0.4), Inches(1.6), Inches(1.2), Inches(0.08))
                accent_bar.fill.solid()
                accent_bar.fill.fore_color.rgb = RGBColor(96, 165, 250)
                accent_bar.line.fill.background()

                header_text_color = RGBColor(255, 255, 255)
                content_text_color = text_color

            elif template_name == "elegant_dark":
                # Dark gradient background
                bg_gradient = [
                    (RGBColor(31, 41, 55), 0, 0, 10, 3.75),    # Gray-800
                    (RGBColor(17, 24, 39), 0, 3.75, 10, 3.75)  # Gray-900
                ]
                for color, x, y, w, h in bg_gradient:
                    rect = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(h))
                    rect.fill.solid()
                    rect.fill.fore_color.rgb = color
                    rect.line.fill.background()

                # Modern dark header card
                header_gradient = [
                    (RGBColor(17, 24, 39), 0.4, 0.4, 4.6, 1.1),  # Gray-900
                    (RGBColor(31, 41, 55), 5, 0.4, 4.6, 1.1)     # Gray-800
                ]
                for color, x, y, w, h in header_gradient:
                    rect = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(h))
                    rect.fill.solid()
                    rect.fill.fore_color.rgb = color
                    rect.line.fill.background()

                # Yellow accent bar under header
                accent_bar = slide.shapes.add_shape(1, Inches(0.4), Inches(1.6), Inches(1.2), Inches(0.08))
                accent_bar.fill.solid()
                accent_bar.fill.fore_color.rgb = RGBColor(251, 191, 36)  # Yellow-400
                accent_bar.line.fill.background()

                header_text_color = RGBColor(255, 255, 255)
                content_text_color = RGBColor(229, 231, 235)  # Gray-200

            elif template_name == "vibrant_education":
                # White background
                bg = slide.shapes.add_shape(1, Inches(0), Inches(1.2), Inches(10), Inches(6.3))
                bg.fill.solid()
                bg.fill.fore_color.rgb = RGBColor(255, 255, 255)
                bg.line.fill.background()

                # Colorful gradient header
                header_colors = [(RGBColor(138, 43, 226), 0), (RGBColor(52, 152, 219), 5)]
                for color, x_offset in header_colors:
                    header = slide.shapes.add_shape(1, Inches(x_offset), Inches(0), Inches(5), Inches(1.2))
                    header.fill.solid()
                    header.fill.fore_color.rgb = color
                    header.line.fill.background()

                # Colorful accent bars
                accent_bars = [(RGBColor(46, 204, 113), 1.2), (RGBColor(255, 140, 0), 3.0),
                               (RGBColor(255, 215, 0), 4.8), (RGBColor(52, 152, 219), 6.6)]
                for color, y_pos in accent_bars:
                    bar = slide.shapes.add_shape(1, Inches(0), Inches(y_pos), Inches(0.15), Inches(1.5))
                    bar.fill.solid()
                    bar.fill.fore_color.rgb = color
                    bar.line.fill.background()

                # Corner decoration
                corner = slide.shapes.add_shape(1, Inches(9), Inches(6.5), Inches(1), Inches(1))
                corner.fill.solid()
                corner.fill.fore_color.rgb = RGBColor(255, 140, 0)
                corner.line.fill.background()
                corner.rotation = 45

                header_text_color = RGBColor(255, 255, 255)
                content_text_color = text_color

            else:
                # Generic content slide for new templates (ocean_breeze, sunset_glow, forest_fresh, royal_purple, corporate_slate)
                # Determine header colors based on template
                if template_name == "ocean_breeze":
                    header_color1 = RGBColor(2, 132, 199)
                    header_color2 = RGBColor(6, 182, 212)
                elif template_name == "sunset_glow":
                    header_color1 = RGBColor(249, 115, 22)
                    header_color2 = RGBColor(236, 72, 153)
                elif template_name == "forest_fresh":
                    header_color1 = RGBColor(6, 95, 70)
                    header_color2 = RGBColor(5, 150, 105)
                elif template_name == "royal_purple":
                    header_color1 = RGBColor(109, 40, 217)
                    header_color2 = RGBColor(139, 92, 246)
                elif template_name == "corporate_slate":
                    header_color1 = RGBColor(30, 41, 59)
                    header_color2 = RGBColor(59, 130, 246)
                else:
                    header_color1 = RGBColor(59, 130, 246)
                    header_color2 = RGBColor(99, 102, 241)

                # Header with gradient effect
                header1 = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(5), Inches(1.3))
                header1.fill.solid()
                header1.fill.fore_color.rgb = header_color1
                header1.line.fill.background()

                header2 = slide.shapes.add_shape(1, Inches(5), Inches(0), Inches(5), Inches(1.3))
                header2.fill.solid()
                header2.fill.fore_color.rgb = header_color2
                header2.line.fill.background()

                # White content area
                content_bg = slide.shapes.add_shape(1, Inches(0), Inches(1.3), Inches(10), Inches(6.2))
                content_bg.fill.solid()
                content_bg.fill.fore_color.rgb = RGBColor(255, 255, 255)
                content_bg.line.fill.background()

                # Accent bar on left
                accent_bar = slide.shapes.add_shape(1, Inches(0), Inches(1.3), Inches(0.25), Inches(6.2))
                accent_bar.fill.solid()
                accent_bar.fill.fore_color.rgb = accent_color
                accent_bar.line.fill.background()

                header_text_color = RGBColor(255, 255, 255)
                content_text_color = text_color

            # Add content card (white/dark card for content area)
            if template_name == "elegant_dark":
                card_color = RGBColor(55, 65, 81)  # Gray-700
            else:
                card_color = RGBColor(255, 255, 255)  # White

            content_card = slide.shapes.add_shape(
                1,  # Rectangle
                Inches(0.4), Inches(1.9), Inches(9.2), Inches(5.4)
            )
            content_card.fill.solid()
            content_card.fill.fore_color.rgb = card_color
            content_card.line.color.rgb = RGBColor(209, 213, 219) if template_name != "elegant_dark" else RGBColor(75, 85, 99)
            content_card.line.width = Pt(1)

            # Add slide title in header
            title_box = slide.shapes.add_textbox(Inches(0.6), Inches(0.6), Inches(8.8), Inches(0.7))
            title_frame = title_box.text_frame
            title_frame.text = slide_data['title']
            title_para = title_frame.paragraphs[0]
            title_para.font.size = Pt(36)
            title_para.font.bold = True
            title_para.font.color.rgb = header_text_color

            # Decide layout: if image needed, split content and image area
            has_image = slide_data.get('has_image', False)

            if has_image:
                # Split slide: left side for content, right side for image
                content_left = Inches(0.7)
                content_width = Inches(4.3)
                image_left = Inches(5.3)
                image_width = Inches(3.9)
                image_top = Inches(2.1)
                image_height = Inches(4.9)
            else:
                # Full width for content
                content_left = Inches(0.7)
                content_width = Inches(8.6)

            # Add bullet points - positioned inside card with proper margins
            content_top = Inches(2.1)
            content_height = Inches(4.9)

            content_box = slide.shapes.add_textbox(content_left, content_top, content_width, content_height)
            text_frame = content_box.text_frame
            text_frame.word_wrap = True
            text_frame.auto_size = None  # Disable auto-size to prevent overflow
            text_frame.margin_left = Inches(0.15)
            text_frame.margin_right = Inches(0.15)
            text_frame.margin_top = Inches(0.1)
            text_frame.margin_bottom = Inches(0.1)

            # Add each bullet point with modern formatting
            for j, point in enumerate(slide_data.get('content', [])):
                if j == 0:
                    # First paragraph already exists
                    para = text_frame.paragraphs[0]
                else:
                    # Add new paragraph for subsequent points
                    para = text_frame.add_paragraph()

                para.level = 0  # Bullet level (0 = main bullet)
                para.space_before = Pt(8)
                para.space_after = Pt(6)
                para.line_spacing = 1.15  # Tighter line height to prevent overflow

                # Parse and apply bold formatting
                # Split by ** markers to find bold text
                import re
                parts = re.split(r'\*\*(.*?)\*\*', point)

                for idx, part in enumerate(parts):
                    if not part:
                        continue
                    run = para.add_run()
                    run.text = part
                    run.font.name = 'Segoe UI'
                    run.font.size = Pt(14)  # Reduced from 16
                    run.font.color.rgb = content_text_color

                    # Bold for parts between ** markers (odd indices)
                    if idx % 2 == 1:
                        run.font.bold = True
                        run.font.size = Pt(16)  # Reduced from 18
                        if template_name != "elegant_dark":
                            run.font.color.rgb = accent_color  # Accent color for bold terms

            # Add image (from Pexels or placeholder)
            if has_image:
                image_added = False

                # Try to fetch and add real image from Pexels
                if PEXELS_AVAILABLE:
                    try:
                        # Get image for this slide
                        image_path = get_image_for_slide(slide_data['title'], i)

                        if image_path and os.path.exists(image_path):
                            # Add the actual image
                            slide.shapes.add_picture(
                                image_path,
                                image_left, image_top,
                                width=image_width, height=image_height
                            )
                            image_added = True
                            print(f"  ✓ Added image from Pexels for slide {i+1}")
                    except Exception as e:
                        print(f"  ⚠ Could not add Pexels image: {e}")

                # Fallback to placeholder if image wasn't added
                if not image_added:
                    # Create gradient background for placeholder
                    placeholder_gradient = [
                        (RGBColor(224, 231, 255), image_left, image_top, image_width, image_height/3),
                        (RGBColor(221, 214, 254), image_left, image_top + image_height/3, image_width, image_height/3),
                        (RGBColor(251, 207, 232), image_left, image_top + 2*image_height/3, image_width, image_height/3)
                    ]
                    for color, x, y, w, h in placeholder_gradient:
                        rect = slide.shapes.add_shape(1, x, y, w, h)
                        rect.fill.solid()
                        rect.fill.fore_color.rgb = color
                        rect.line.fill.background()

                    # Add dashed border
                    border = slide.shapes.add_shape(1, image_left, image_top, image_width, image_height)
                    border.fill.background()
                    border.line.color.rgb = RGBColor(99, 102, 241)
                    border.line.width = Pt(3)
                    border.line.dash_style = 2

                    # Add icon and text in center
                    icon_box = slide.shapes.add_textbox(
                        image_left + Inches(0.5),
                        image_top + image_height/2 - Inches(0.4),
                        image_width - Inches(1),
                        Inches(0.8)
                    )
                    icon_frame = icon_box.text_frame
                    icon_frame.text = "📸 Visual Content\nAdd image or diagram"
                    icon_para = icon_frame.paragraphs[0]
                    icon_para.font.size = Pt(16)
                    icon_para.font.bold = True
                    icon_para.font.color.rgb = RGBColor(79, 70, 229)
                    icon_para.alignment = PP_ALIGN.CENTER
                    icon_para.line_spacing = 1.4

        # Add speaker notes to the slide
        notes_slide = slide.notes_slide
        notes_frame = notes_slide.notes_text_frame
        notes_frame.text = slide_data.get('notes', '')

    # Create output directory if it doesn't exist
    os.makedirs('output', exist_ok=True)

    # Save the presentation
    output_path = f"output/{output_filename}"
    prs.save(output_path)

    print(f"\nPresentation saved successfully: {output_path}")
    return output_path


def generate_ppt(topic, template_name="ocean_breeze", num_slides=10, output_filename=None):
    """
    Main function that orchestrates the entire PPT generation process.
    This is the function you'll call from the Streamlit app.

    Args:
        topic: The topic/prompt for the presentation
        template_name: Which template to use (ocean_breeze, sunset_glow, forest_fresh, etc.)
        num_slides: Number of slides to generate (default: 10 for educational content)
        output_filename: Name for output file (auto-generated if None)

    Returns:
        Path to the generated PowerPoint file
    """

    # Auto-generate filename if not provided
    if output_filename is None:
        # Create a safe filename from the topic
        safe_topic = "".join(c if c.isalnum() else "_" for c in topic)
        safe_topic = safe_topic[:50]  # Limit length
        output_filename = f"{safe_topic}_presentation.pptx"

    # Step 1: Initialize OpenAI client
    print("Step 1: Initializing OpenAI client...")
    client = initialize_openai_client()

    # Step 2: Generate slide content using AI
    print(f"\nStep 2: Generating content for {num_slides} slides on topic: '{topic}'")
    slides_data = generate_slide_content(client, topic, num_slides)

    # Step 3: Create the PowerPoint presentation
    print(f"\nStep 3: Creating PowerPoint with template '{template_name}'")
    output_path = create_presentation_from_template(template_name, slides_data, output_filename)

    print("\n[SUCCESS] AI-powered presentation generated successfully!")
    return output_path


# ========== TEST FUNCTION (optional) ==========
# You can run this file directly to test the generation
if __name__ == "__main__":
    print("AI PPT Generator - Test Mode")
    print("=" * 50)

    # Test parameters
    test_topic = "Artificial Intelligence in Education"
    test_template = "modern_blue"  # or "elegant_dark"
    test_slides = 5

    try:
        result_path = generate_ppt(
            topic=test_topic,
            template_name=test_template,
            num_slides=test_slides
        )
        print(f"\nTest successful! Check your presentation at: {result_path}")
    except Exception as e:
        print(f"\nError during generation: {e}")
        import traceback
        traceback.print_exc()
