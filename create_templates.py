"""
This script creates professional PowerPoint templates for the AI PPT Maker.
Run this once to generate template files in the templates folder.
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import os

# Create templates directory if it doesn't exist
os.makedirs('templates', exist_ok=True)

# ========== TEMPLATE 1: MODERN BLUE (Corporate Professional) ==========
print("Creating Template 1: Modern Blue...")

prs1 = Presentation()
prs1.slide_width = Inches(10)
prs1.slide_height = Inches(7.5)

# Define the Modern Blue color scheme
BLUE_DARK = RGBColor(26, 35, 126)      # Dark blue for headers
BLUE_ACCENT = RGBColor(33, 150, 243)   # Bright blue for accents
WHITE = RGBColor(255, 255, 255)        # White for text
GRAY_LIGHT = RGBColor(245, 245, 245)   # Light gray for backgrounds

# Add Title Slide Layout
title_slide = prs1.slides.add_slide(prs1.slide_layouts[6])  # Blank layout

# Title slide background (gradient effect with top blue bar)
top_bar = title_slide.shapes.add_shape(
    1,  # Rectangle
    Inches(0), Inches(0),
    Inches(10), Inches(2)
)
top_bar.fill.solid()
top_bar.fill.fore_color.rgb = BLUE_DARK
top_bar.line.fill.background()

# Bottom section background
bottom_section = title_slide.shapes.add_shape(
    1,  # Rectangle
    Inches(0), Inches(2),
    Inches(10), Inches(5.5)
)
bottom_section.fill.solid()
bottom_section.fill.fore_color.rgb = GRAY_LIGHT
bottom_section.line.fill.background()

# Add accent line
accent_line = title_slide.shapes.add_shape(
    1,  # Rectangle
    Inches(0), Inches(2),
    Inches(10), Inches(0.1)
)
accent_line.fill.solid()
accent_line.fill.fore_color.rgb = BLUE_ACCENT
accent_line.line.fill.background()

# Add Content Slide Layout
content_slide = prs1.slides.add_slide(prs1.slide_layouts[6])

# Content slide header bar
header_bar = content_slide.shapes.add_shape(
    1,
    Inches(0), Inches(0),
    Inches(10), Inches(1)
)
header_bar.fill.solid()
header_bar.fill.fore_color.rgb = BLUE_DARK
header_bar.line.fill.background()

# Content area background
content_bg = content_slide.shapes.add_shape(
    1,
    Inches(0), Inches(1),
    Inches(10), Inches(6.5)
)
content_bg.fill.solid()
content_bg.fill.fore_color.rgb = WHITE
content_bg.line.fill.background()

# Add accent sidebar
sidebar = content_slide.shapes.add_shape(
    1,
    Inches(0), Inches(1),
    Inches(0.2), Inches(6.5)
)
sidebar.fill.solid()
sidebar.fill.fore_color.rgb = BLUE_ACCENT
sidebar.line.fill.background()

# Save Template 1
prs1.save('templates/modern_blue.pptx')
print("[OK] Template 1 saved: templates/modern_blue.pptx")


# ========== TEMPLATE 2: ELEGANT DARK (Modern Premium) ==========
print("\nCreating Template 2: Elegant Dark...")

prs2 = Presentation()
prs2.slide_width = Inches(10)
prs2.slide_height = Inches(7.5)

# Define the Elegant Dark color scheme
DARK_BG = RGBColor(18, 18, 18)         # Almost black background
GOLD_ACCENT = RGBColor(255, 193, 7)    # Gold accent color
DARK_GRAY = RGBColor(33, 33, 33)       # Dark gray for sections
TEXT_WHITE = RGBColor(250, 250, 250)   # Off-white for text

# Add Title Slide Layout
title_slide2 = prs2.slides.add_slide(prs2.slide_layouts[6])

# Dark background
bg = title_slide2.shapes.add_shape(
    1,
    Inches(0), Inches(0),
    Inches(10), Inches(7.5)
)
bg.fill.solid()
bg.fill.fore_color.rgb = DARK_BG
bg.line.fill.background()

# Gold accent strip (diagonal design element)
accent1 = title_slide2.shapes.add_shape(
    1,
    Inches(8), Inches(0),
    Inches(2), Inches(7.5)
)
accent1.fill.solid()
accent1.fill.fore_color.rgb = GOLD_ACCENT
accent1.line.fill.background()
accent1.rotation = 15  # Slight tilt for modern look

# Add Content Slide Layout
content_slide2 = prs2.slides.add_slide(prs2.slide_layouts[6])

# Dark background
bg2 = content_slide2.shapes.add_shape(
    1,
    Inches(0), Inches(0),
    Inches(10), Inches(7.5)
)
bg2.fill.solid()
bg2.fill.fore_color.rgb = DARK_BG
bg2.line.fill.background()

# Header section with dark gray
header_section = content_slide2.shapes.add_shape(
    1,
    Inches(0), Inches(0),
    Inches(10), Inches(1.2)
)
header_section.fill.solid()
header_section.fill.fore_color.rgb = DARK_GRAY
header_section.line.fill.background()

# Gold accent line under header
gold_line = content_slide2.shapes.add_shape(
    1,
    Inches(0.5), Inches(1.1),
    Inches(2), Inches(0.08)
)
gold_line.fill.solid()
gold_line.fill.fore_color.rgb = GOLD_ACCENT
gold_line.line.fill.background()

# Save Template 2
prs2.save('templates/elegant_dark.pptx')
print("[OK] Template 2 saved: templates/elegant_dark.pptx")

print("\n[SUCCESS] All templates created successfully!")
print("\nTemplates available:")
print("  1. modern_blue.pptx - Professional corporate design with blue theme")
print("  2. elegant_dark.pptx - Premium modern design with dark theme")
