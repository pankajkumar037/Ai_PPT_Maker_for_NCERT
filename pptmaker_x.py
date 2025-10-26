"""
PPTMaker-X: Single Conversational Agent with Template-Based Design
Uses pre-designed templates and allows interactive editing of individual slides
"""

import os
import json
import random
from typing import Dict, List, Any, Optional
from dotenv import load_dotenv
from openai import OpenAI
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from PIL import Image
import io

# Load environment variables
load_dotenv()


class TemplatePPTGenerator:
    """
    Template-based PowerPoint generator
    Uses pre-designed templates for consistent, professional layouts
    """

    # Available themes matching our clean templates
    THEMES = {
        'modern_blue': {'name': 'Modern Blue', 'file': 'templates/modern_blue_clean.pptx'},
        'elegant_purple': {'name': 'Elegant Purple', 'file': 'templates/elegant_purple_clean.pptx'},
        'corporate_green': {'name': 'Corporate Green', 'file': 'templates/corporate_green_clean.pptx'},
        'vibrant_orange': {'name': 'Vibrant Orange', 'file': 'templates/vibrant_orange_clean.pptx'},
        'dark_professional': {'name': 'Dark Professional', 'file': 'templates/dark_professional_clean.pptx'},
        'minimal_gray': {'name': 'Minimal Gray', 'file': 'templates/minimal_gray_clean.pptx'},
        'ocean_teal': {'name': 'Ocean Teal', 'file': 'templates/ocean_teal_clean.pptx'},
        'sunset_red': {'name': 'Sunset Red', 'file': 'templates/sunset_red_clean.pptx'},
        'royal_indigo': {'name': 'Royal Indigo', 'file': 'templates/royal_indigo_clean.pptx'},
        'forest_brown': {'name': 'Forest Brown', 'file': 'templates/forest_brown_clean.pptx'}
    }

    # Layout mapping (which slide number in template to use for which purpose)
    LAYOUT_MAP = {
        'title': 0,           # Slide 1: Hero title
        'section': 1,         # Slide 2: Section divider
        'content': 2,         # Slide 3: Content + visual
        'features': 3,        # Slide 4: Feature cards
        'data': 4,            # Slide 5: Data/chart
        'quote': 5,           # Slide 6: Quote
        'split': 6,           # Slide 7: Split screen
        'timeline': 7,        # Slide 8: Timeline
        'stat': 8,            # Slide 9: Big number
        'hero_image': 9,      # Slide 10: Image hero
        'comparison': 10,     # Slide 11: Comparison
        'team': 11,           # Slide 12: Team grid
        'bullets': 12,        # Slide 13: Bullet list
        'cta': 13,            # Slide 14: Call to action
        'dashboard': 14       # Slide 15: Dashboard
    }

    def __init__(self, theme_key: str = 'modern_blue'):
        """Initialize generator with a specific theme"""
        self.theme_key = theme_key
        self.theme = self.THEMES.get(theme_key, self.THEMES['modern_blue'])
        self.template_path = self.theme['file']

        # Load template presentation - fallback to any available template
        if os.path.exists(self.template_path):
            self.template_prs = Presentation(self.template_path)
            print(f"Using template: {self.template_path}")
        else:
            print(f"Warning: Template {self.template_path} not found. Trying alternatives...")
            # Try to find ANY available template
            for alt_theme_key, alt_theme in self.THEMES.items():
                alt_path = alt_theme['file']
                if os.path.exists(alt_path):
                    self.template_path = alt_path
                    self.template_prs = Presentation(alt_path)
                    self.theme_key = alt_theme_key
                    print(f"Using alternative template: {alt_path}")
                    break
            else:
                # No templates found at all
                print("No templates found! Using blank presentation.")
                self.template_prs = Presentation()

        # Create new presentation for output
        self.prs = Presentation()
        self.prs.slide_width = Inches(10)
        self.prs.slide_height = Inches(5.625)

    def copy_slide(self, template_slide_idx: int, content: Dict) -> None:
        """
        Copy a slide from template and populate with content

        Args:
            template_slide_idx: Index of slide in template to copy
            content: Dictionary with text content to populate
        """
        # Get template slide
        if template_slide_idx >= len(self.template_prs.slides):
            print(f"Warning: Template slide {template_slide_idx} not found")
            return

        template_slide = self.template_prs.slides[template_slide_idx]

        # Create new slide from blank layout
        new_slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])

        # Copy all shapes from template
        for shape in template_slide.shapes:
            # Copy the shape's element
            el = shape.element
            new_el = self._copy_element(el)
            new_slide.shapes._spTree.insert_element_before(new_el, 'p:extLst')

        # Update text content based on content dict
        self._populate_slide_content(new_slide, content)

    def _copy_element(self, element):
        """Deep copy an XML element"""
        from copy import deepcopy
        return deepcopy(element)

    def _populate_slide_content(self, slide, content: Dict) -> None:
        """
        Populate slide with actual content - AGGRESSIVE VERSION

        Args:
            slide: PowerPoint slide object
            content: Dict with keys like 'title', 'subtitle', 'bullets', 'text'
        """
        # Collect all text shapes sorted by top position (top to bottom)
        text_shapes = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                text_shapes.append(shape)

        # Sort by vertical position (top to bottom)
        text_shapes.sort(key=lambda s: s.top)

        # Strategy: Replace text boxes from top to bottom with our content
        shape_idx = 0

        # First shape usually = title
        if shape_idx < len(text_shapes) and 'title' in content:
            text_frame = text_shapes[shape_idx].text_frame
            text_frame.clear()
            p = text_frame.paragraphs[0]
            p.text = content['title']
            p.font.bold = True
            p.font.size = Pt(40)
            for run in p.runs:
                run.font.bold = True
                run.font.size = Pt(40)
            shape_idx += 1

        # Second shape usually = subtitle or content
        if shape_idx < len(text_shapes):
            text_frame = text_shapes[shape_idx].text_frame

            # If we have subtitle, use it
            if 'subtitle' in content:
                text_frame.clear()
                p = text_frame.paragraphs[0]
                p.text = content['subtitle']
                p.font.size = Pt(24)
                for run in p.runs:
                    run.font.size = Pt(24)
                shape_idx += 1

            # If we have bullets, populate them
            elif 'bullets' in content:
                text_frame.clear()
                for i, bullet in enumerate(content['bullets'][:6]):  # Max 6 bullets
                    if i == 0:
                        p = text_frame.paragraphs[0]
                    else:
                        p = text_frame.add_paragraph()
                    p.text = bullet
                    p.level = 0
                    p.font.size = Pt(20)
                    for run in p.runs:
                        run.font.size = Pt(20)
                        run.font.bold = False
                shape_idx += 1

            # If we have text content
            elif 'text' in content:
                text_frame.clear()
                p = text_frame.paragraphs[0]
                p.text = content['text']
                p.font.size = Pt(20)
                for run in p.runs:
                    run.font.size = Pt(20)
                shape_idx += 1

            # If we have stat
            elif 'stat' in content:
                text_frame.clear()
                p = text_frame.paragraphs[0]
                p.text = content['stat']
                p.font.bold = True
                p.font.size = Pt(72)
                for run in p.runs:
                    run.font.bold = True
                    run.font.size = Pt(72)
                shape_idx += 1

        # Third shape (if exists) - for additional content
        if shape_idx < len(text_shapes):
            text_frame = text_shapes[shape_idx].text_frame

            # If we haven't used bullets yet and we have them
            if 'bullets' in content and 'subtitle' not in content:
                # Already used above
                pass

            # If we have description for stat
            elif 'description' in content and 'stat' in content:
                text_frame.clear()
                p = text_frame.paragraphs[0]
                p.text = content['description']
                p.font.size = Pt(20)
                for run in p.runs:
                    run.font.size = Pt(20)
                shape_idx += 1

            # If we have context
            elif 'context' in content:
                text_frame.clear()
                p = text_frame.paragraphs[0]
                p.text = content['context']
                p.font.size = Pt(18)
                for run in p.runs:
                    run.font.size = Pt(18)
                shape_idx += 1

    def save(self, output_path: str) -> str:
        """Save the presentation"""
        os.makedirs(os.path.dirname(output_path) if os.path.dirname(output_path) else ".", exist_ok=True)
        self.prs.save(output_path)
        return output_path


class ConversationalPPTAgent:
    """
    Single AI Agent for conversational PowerPoint creation
    Handles outline generation, content creation, and iterative editing
    """

    def __init__(self):
        """Initialize the conversational agent"""
        self.client = OpenAI(api_key=os.getenv("OPENAI_API_KEY"))
        self.model = "gpt-4o-mini"  # Using GPT-4 for better quality

        # Session state
        self.outline = None
        self.theme = None
        self.slides_content = []
        self.ppt_generator = None

    def generate_outline(self, topic: str, num_slides: int) -> Dict:
        """
        STEP 1: Generate presentation outline for user review

        Args:
            topic: Presentation topic
            num_slides: Number of slides requested

        Returns:
            Dict with outline structure
        """
        print("\n" + "="*70)
        print("STEP 1: Generating Presentation Outline")
        print("="*70)

        prompt = f"""Create a presentation outline for: "{topic}"

Requirements:
- Total slides: {num_slides}
- First slide: Title slide
- Last slide: Summary/Conclusion
- Middle slides: Core content covering the topic comprehensively

For each slide, specify:
1. Slide number
2. Slide type (title/content/stat/comparison/etc)
3. Main topic/title
4. Brief description of what it covers

Respond with JSON:
{{
    "presentation_title": "...",
    "slides": [
        {{"number": 1, "type": "title", "topic": "...", "description": "..."}},
        {{"number": 2, "type": "content", "topic": "...", "description": "..."}},
        ...
    ]
}}"""

        response = self.client.chat.completions.create(
            model=self.model,
            messages=[{"role": "user", "content": prompt}],
            temperature=0.7
        )

        outline = self._extract_json(response.choices[0].message.content)
        self.outline = outline

        # Display outline
        print(f"\nPresentation: {outline.get('presentation_title', 'Untitled')}")
        print(f"Total Slides: {len(outline.get('slides', []))}\n")

        for slide in outline.get('slides', []):
            print(f"  Slide {slide['number']}: {slide['topic']}")
            print(f"    Type: {slide['type']}")
            print(f"    Content: {slide['description']}\n")

        return outline

    def select_theme(self, user_preference: Optional[str] = None) -> str:
        """
        STEP 2: Select theme based on presentation topic

        Args:
            user_preference: Optional theme preference from user

        Returns:
            Selected theme key
        """
        print("\n" + "="*70)
        print("STEP 2: Selecting Presentation Theme")
        print("="*70)

        if user_preference and user_preference in TemplatePPTGenerator.THEMES:
            theme_key = user_preference
            print(f"\nUsing user-selected theme: {TemplatePPTGenerator.THEMES[theme_key]['name']}")
        else:
            # AI selects theme based on topic
            if not self.outline:
                theme_key = 'modern_blue'
            else:
                prompt = f"""Based on this presentation topic: "{self.outline.get('presentation_title', '')}"

Select the most appropriate theme from:
{json.dumps(list(TemplatePPTGenerator.THEMES.keys()))}

Consider:
- modern_blue: Technology, business, professional
- elegant_purple: Creative, luxury, artistic
- corporate_green: Sustainability, finance, growth
- vibrant_orange: Energy, innovation, dynamic
- dark_professional: Tech, modern, sophisticated
- minimal_gray: Minimalist, clean, neutral
- ocean_teal: Healthcare, wellness, calm
- sunset_red: Passion, urgency, bold
- royal_indigo: Authority, traditional, corporate
- forest_brown: Natural, organic, earthy

Respond with just the theme key (e.g., "modern_blue")"""

                response = self.client.chat.completions.create(
                    model=self.model,
                    messages=[{"role": "user", "content": prompt}],
                    temperature=0.3
                )

                theme_key = response.choices[0].message.content.strip().strip('"')
                if theme_key not in TemplatePPTGenerator.THEMES:
                    theme_key = 'modern_blue'

            print(f"\nAI selected theme: {TemplatePPTGenerator.THEMES[theme_key]['name']}")

        self.theme = theme_key
        self.ppt_generator = TemplatePPTGenerator(theme_key)

        return theme_key

    def generate_content(self) -> List[Dict]:
        """
        STEP 3: Generate detailed content for each slide

        Returns:
            List of slide content dictionaries
        """
        print("\n" + "="*70)
        print("STEP 3: Generating Slide Content")
        print("="*70)

        slides_content = []

        for slide_info in self.outline.get('slides', []):
            print(f"\n  Generating content for Slide {slide_info['number']}: {slide_info['topic']}...")

            # Generate content based on slide type
            content = self._generate_slide_content(slide_info)
            slides_content.append(content)

        self.slides_content = slides_content
        print(f"\n  Generated content for {len(slides_content)} slides")

        return slides_content

    def _generate_slide_content(self, slide_info: Dict) -> Dict:
        """
        Generate detailed content for a specific slide

        Args:
            slide_info: Slide information from outline

        Returns:
            Dict with detailed content (title, bullets, text, etc.)
        """
        slide_type = slide_info['type']

        # Content generation prompt based on slide type
        if slide_type == 'title':
            prompt = f"""Create content for a title slide:
Topic: {slide_info['topic']}
Description: {slide_info['description']}

Respond with JSON:
{{"title": "Main Title", "subtitle": "Engaging subtitle"}}"""

        elif slide_type in ['content', 'bullets']:
            prompt = f"""Create content for a content slide:
Topic: {slide_info['topic']}
Description: {slide_info['description']}

Respond with JSON:
{{"title": "Slide Title", "bullets": ["Point 1 with details", "Point 2 with details", "Point 3 with details", "Point 4 with details"]}}"""

        elif slide_type == 'stat':
            prompt = f"""Create content for a statistics/big number slide:
Topic: {slide_info['topic']}
Description: {slide_info['description']}

Respond with JSON:
{{"stat": "95%", "description": "What this statistic means", "context": "Additional context"}}"""

        else:
            # Default content structure
            prompt = f"""Create content for a slide:
Type: {slide_type}
Topic: {slide_info['topic']}
Description: {slide_info['description']}

Respond with JSON:
{{"title": "Slide Title", "text": "Main content text for this slide"}}"""

        response = self.client.chat.completions.create(
            model=self.model,
            messages=[{"role": "user", "content": prompt}],
            temperature=0.7
        )

        content = self._extract_json(response.choices[0].message.content)
        content['slide_number'] = slide_info['number']
        content['slide_type'] = slide_type

        return content

    def apply_user_feedback(self, feedback: str) -> None:
        """
        STEP 4: Apply user feedback to modify slides

        Args:
            feedback: User's feedback (e.g., "Make slide 3 more colorful", "Add stats to slide 5")
        """
        print("\n" + "="*70)
        print("STEP 4: Applying User Feedback")
        print("="*70)
        print(f"\nFeedback: {feedback}\n")

        # AI analyzes feedback and determines what to change
        prompt = f"""Analyze this user feedback about their presentation:
"{feedback}"

Current presentation has {len(self.slides_content)} slides.
Slide topics: {[s.get('title', s.get('topic', '')) for s in self.slides_content]}

Determine:
1. Which slide(s) need to be modified (slide numbers)
2. What specific changes to make
3. If it's a content change, theme change, or layout change

Respond with JSON:
{{
    "action": "modify_slide/change_theme/add_slide/global_change",
    "target_slides": [slide numbers to modify],
    "changes": "Detailed description of changes to make",
    "new_content": {{"title": "...", "bullets": [...]}} (if modifying content)
}}"""

        response = self.client.chat.completions.create(
            model=self.model,
            messages=[{"role": "user", "content": prompt}],
            temperature=0.5
        )

        action_plan = self._extract_json(response.choices[0].message.content)

        print(f"  Action: {action_plan.get('action', 'unknown')}")
        print(f"  Target slides: {action_plan.get('target_slides', [])}")
        print(f"  Changes: {action_plan.get('changes', 'None specified')}")

        # Apply the changes
        if action_plan.get('action') == 'modify_slide':
            for slide_num in action_plan.get('target_slides', []):
                if 0 < slide_num <= len(self.slides_content):
                    # Update the specific slide content
                    if 'new_content' in action_plan:
                        self.slides_content[slide_num - 1].update(action_plan['new_content'])
                    print(f"  Updated slide {slide_num}")

        elif action_plan.get('action') == 'change_theme':
            # Change presentation theme
            new_theme = action_plan.get('new_theme', 'modern_blue')
            self.select_theme(new_theme)

        print("\n  Feedback applied successfully!")

    def self_review(self) -> Dict:
        """
        STEP 5: Agent reviews the presentation and suggests improvements

        Returns:
            Dict with review results and suggestions
        """
        print("\n" + "="*70)
        print("STEP 5: Self-Review & Quality Check")
        print("="*70)

        # AI reviews the complete presentation
        prompt = f"""Review this presentation for quality:

Title: {self.outline.get('presentation_title', '')}
Theme: {self.theme}
Number of slides: {len(self.slides_content)}

Content summary:
{json.dumps(self.slides_content, indent=2)}

Check for:
1. Content quality and relevance
2. Flow and logical progression
3. Completeness of coverage
4. Balance of information
5. Any missing important points

Respond with JSON:
{{
    "overall_quality": "excellent/good/needs_improvement",
    "score": 85,
    "strengths": ["What works well"],
    "suggestions": ["Specific improvements"],
    "missing_topics": ["Topics that should be added"]
}}"""

        response = self.client.chat.completions.create(
            model=self.model,
            messages=[{"role": "user", "content": prompt}],
            temperature=0.3
        )

        review = self._extract_json(response.choices[0].message.content)

        print(f"\n  Quality: {review.get('overall_quality', 'N/A')}")
        print(f"  Score: {review.get('score', 0)}/100")
        print(f"\n  Strengths:")
        for strength in review.get('strengths', []):
            print(f"    + {strength}")

        print(f"\n  Suggestions:")
        for suggestion in review.get('suggestions', []):
            print(f"    - {suggestion}")

        return review

    def build_single_slide(self, slide_index: int, output_path: str) -> str:
        """
        Build a single slide and save presentation (Gamma.AI style)

        Args:
            slide_index: Index of slide to build (0-based)
            output_path: Where to save the presentation

        Returns:
            Path to saved presentation
        """
        if not self.ppt_generator:
            print("Error: No theme selected")
            return None

        if slide_index >= len(self.slides_content):
            print(f"Error: Slide index {slide_index} out of range")
            return None

        slide_content = self.slides_content[slide_index]
        slide_type = slide_content.get('slide_type', 'content')

        # Map slide type to template layout
        layout_key = self._map_slide_type_to_layout(slide_type)
        template_idx = TemplatePPTGenerator.LAYOUT_MAP.get(layout_key, 2)

        # Copy template slide and populate with content
        self.ppt_generator.copy_slide(template_idx, slide_content)
        print(f"  Created slide {slide_content.get('slide_number', '?')}: {slide_content.get('title', 'Untitled')}")

        # Save presentation
        saved_path = self.ppt_generator.save(output_path)

        return saved_path

    def modify_slide(self, slide_index: int, feedback: str, ppt_path: str) -> Dict:
        """
        Modify a specific slide IN-PLACE based on user feedback
        (AutoGen agent modifies the PPTX slide directly)

        Args:
            slide_index: Index of slide to modify
            feedback: User feedback for this specific slide
            ppt_path: Path to PPTX file to modify

        Returns:
            Updated slide content
        """
        if slide_index >= len(self.slides_content):
            return None

        old_content = self.slides_content[slide_index]
        slide_info = self.outline['slides'][slide_index]

        # AutoGen agent interprets feedback and modifies content
        prompt = f"""You are an AutoGen conversational agent. Modify this slide content based on user feedback.

Current slide content:
{json.dumps(old_content, indent=2)}

User feedback: "{feedback}"

Apply the changes directly. Return ONLY the modified fields as JSON.
If user wants to change title, return {{"title": "new title"}}.
If user wants to add bullets, return {{"bullets": ["point 1", "point 2"]}}.
Return only what changed."""

        response = self.client.chat.completions.create(
            model=self.model,
            messages=[{"role": "user", "content": prompt}],
            temperature=0.5
        )

        changes = self._extract_json(response.choices[0].message.content)

        # Apply changes to existing content (in-place modification)
        old_content.update(changes)

        # Update in slides_content
        self.slides_content[slide_index] = old_content

        # Rebuild ONLY this slide in the PPTX file
        prs = Presentation(ppt_path)
        if slide_index < len(prs.slides):
            # Get the slide
            slide = prs.slides[slide_index]

            # Modify the slide directly with updated content
            self.ppt_generator._populate_slide_content(slide, old_content)

            # Save the PPTX file
            prs.save(ppt_path)
            print(f"  Modified slide {slide_index + 1} in-place in {ppt_path}")

        return old_content

    def build_presentation(self, output_path: str) -> str:
        """
        STEP 6: Build the final PowerPoint file

        Args:
            output_path: Where to save the presentation

        Returns:
            Path to saved presentation
        """
        print("\n" + "="*70)
        print("STEP 6: Building PowerPoint Presentation")
        print("="*70)

        if not self.ppt_generator:
            print("Error: No theme selected")
            return None

        # Create slides using template
        for slide_content in self.slides_content:
            slide_type = slide_content.get('slide_type', 'content')

            # Map slide type to template layout
            layout_key = self._map_slide_type_to_layout(slide_type)
            template_idx = TemplatePPTGenerator.LAYOUT_MAP.get(layout_key, 2)

            # Copy template slide and populate with content
            self.ppt_generator.copy_slide(template_idx, slide_content)
            print(f"  Created slide {slide_content.get('slide_number', '?')}: {slide_content.get('title', 'Untitled')}")

        # Save presentation
        saved_path = self.ppt_generator.save(output_path)

        print(f"\n  Presentation saved to: {saved_path}")
        print(f"  Total slides: {len(self.ppt_generator.prs.slides)}")

        return saved_path

    def _map_slide_type_to_layout(self, slide_type: str) -> str:
        """Map slide type from outline to template layout key"""
        mapping = {
            'title': 'title',
            'content': 'content',
            'bullets': 'bullets',
            'stat': 'stat',
            'statistics': 'stat',
            'data': 'data',
            'comparison': 'comparison',
            'timeline': 'timeline',
            'quote': 'quote',
            'team': 'team',
            'features': 'features',
            'cta': 'cta',
            'summary': 'bullets',
            'conclusion': 'bullets'
        }
        return mapping.get(slide_type, 'content')

    def _extract_json(self, text: str) -> Any:
        """Extract JSON from AI response"""
        if not text:
            return {}

        # Remove markdown code blocks
        if "```json" in text:
            text = text.split("```json")[1].split("```")[0].strip()
        elif "```" in text:
            text = text.split("```")[1].split("```")[0].strip()

        try:
            return json.loads(text)
        except json.JSONDecodeError as e:
            print(f"JSON parsing error: {e}")
            return {} if "{" in text else []


def generate_presentation_interactive(topic: str, num_slides: int = 7) -> str:
    """
    Main function to generate presentation with conversational agent

    Args:
        topic: Presentation topic
        num_slides: Number of slides

    Returns:
        Path to generated presentation
    """
    print("\n" + "="*80)
    print("PPTMaker-X: Conversational AI Agent")
    print("Template-Based Design with Interactive Editing")
    print("="*80)

    # Initialize agent
    agent = ConversationalPPTAgent()

    # Step 1: Generate outline
    outline = agent.generate_outline(topic, num_slides)

    # Step 2: Select theme
    theme = agent.select_theme()

    # Step 3: Generate content
    content = agent.generate_content()

    # Step 4: Self-review
    review = agent.self_review()

    # Step 5: Build presentation
    import datetime
    timestamp = datetime.datetime.now().strftime("%Y%m%d_%H%M%S")
    safe_topic = "".join(c if c.isalnum() else "_" for c in topic[:30])
    output_path = f"output/{safe_topic}_{timestamp}.pptx"

    final_path = agent.build_presentation(output_path)

    print("\n" + "="*80)
    print("Presentation Complete!")
    print("="*80)
    print(f"File: {final_path}")
    print(f"Theme: {TemplatePPTGenerator.THEMES[theme]['name']}")
    print(f"Slides: {len(content)}")
    print("="*80 + "\n")

    return final_path


if __name__ == "__main__":
    # Test the conversational agent
    test_topic = "Renewable Energy and Climate Change"
    result = generate_presentation_interactive(test_topic, num_slides=6)
    print(f"\nTest complete! Presentation saved to: {result}")
