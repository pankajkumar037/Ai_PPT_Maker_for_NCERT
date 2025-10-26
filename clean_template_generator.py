"""
Clean Professional PowerPoint Template Generator
10 Different Professional Templates WITH proper content (NO weird placeholders)
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
import os

class CleanTemplateGenerator:
    """Generates clean professional templates with proper text content"""

    THEMES = {
        'modern_blue': {
            'name': 'Modern Blue',
            'primary': (41, 98, 255),
            'secondary': (21, 101, 192),
            'accent': (0, 188, 212),
            'bg': (255, 255, 255),
            'bg_alt': (245, 247, 250),
            'text': (33, 33, 33),
            'text_light': (117, 117, 117),
            'heading': 'Calibri',
            'body': 'Calibri'
        },
        'elegant_purple': {
            'name': 'Elegant Purple',
            'primary': (103, 58, 183),
            'secondary': (142, 36, 170),
            'accent': (255, 64, 129),
            'bg': (255, 255, 255),
            'bg_alt': (248, 245, 250),
            'text': (33, 33, 33),
            'text_light': (117, 117, 117),
            'heading': 'Georgia',
            'body': 'Georgia'
        },
        'corporate_green': {
            'name': 'Corporate Green',
            'primary': (27, 94, 32),
            'secondary': (56, 142, 60),
            'accent': (255, 193, 7),
            'bg': (255, 255, 255),
            'bg_alt': (245, 248, 245),
            'text': (33, 33, 33),
            'text_light': (97, 97, 97),
            'heading': 'Arial',
            'body': 'Arial'
        },
        'vibrant_orange': {
            'name': 'Vibrant Orange',
            'primary': (230, 81, 0),
            'secondary': (251, 140, 0),
            'accent': (255, 87, 34),
            'bg': (255, 255, 255),
            'bg_alt': (255, 248, 245),
            'text': (33, 33, 33),
            'text_light': (117, 117, 117),
            'heading': 'Verdana',
            'body': 'Verdana'
        },
        'dark_professional': {
            'name': 'Dark Professional',
            'primary': (224, 224, 224),
            'secondary': (189, 189, 189),
            'accent': (79, 195, 247),
            'bg': (33, 33, 33),
            'bg_alt': (48, 48, 48),
            'text': (245, 245, 245),
            'text_light': (189, 189, 189),
            'heading': 'Segoe UI',
            'body': 'Segoe UI'
        },
        'minimal_gray': {
            'name': 'Minimal Gray',
            'primary': (66, 66, 66),
            'secondary': (97, 97, 97),
            'accent': (255, 152, 0),
            'bg': (255, 255, 255),
            'bg_alt': (250, 250, 250),
            'text': (33, 33, 33),
            'text_light': (117, 117, 117),
            'heading': 'Helvetica',
            'body': 'Helvetica'
        },
        'ocean_teal': {
            'name': 'Ocean Teal',
            'primary': (0, 121, 107),
            'secondary': (0, 150, 136),
            'accent': (255, 235, 59),
            'bg': (255, 255, 255),
            'bg_alt': (224, 247, 250),
            'text': (33, 33, 33),
            'text_light': (117, 117, 117),
            'heading': 'Trebuchet MS',
            'body': 'Trebuchet MS'
        },
        'sunset_red': {
            'name': 'Sunset Red',
            'primary': (183, 28, 28),
            'secondary': (211, 47, 47),
            'accent': (255, 193, 7),
            'bg': (255, 255, 255),
            'bg_alt': (255, 245, 245),
            'text': (33, 33, 33),
            'text_light': (117, 117, 117),
            'heading': 'Garamond',
            'body': 'Garamond'
        },
        'royal_indigo': {
            'name': 'Royal Indigo',
            'primary': (26, 35, 126),
            'secondary': (48, 63, 159),
            'accent': (255, 215, 64),
            'bg': (255, 255, 255),
            'bg_alt': (232, 234, 246),
            'text': (33, 33, 33),
            'text_light': (117, 117, 117),
            'heading': 'Cambria',
            'body': 'Cambria'
        },
        'forest_brown': {
            'name': 'Forest Brown',
            'primary': (62, 39, 35),
            'secondary': (93, 64, 55),
            'accent': (139, 195, 74),
            'bg': (255, 255, 255),
            'bg_alt': (245, 245, 240),
            'text': (33, 33, 33),
            'text_light': (117, 117, 117),
            'heading': 'Times New Roman',
            'body': 'Times New Roman'
        }
    }

    def __init__(self, theme_key):
        self.theme = self.THEMES[theme_key]
        self.theme_key = theme_key
        self.prs = Presentation()
        self.prs.slide_width = Inches(10)
        self.prs.slide_height = Inches(5.625)

    def _rgb(self, key):
        return RGBColor(*self.theme[key])

    def _add_bg(self, slide, color_key='bg'):
        bg = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0), Inches(0), Inches(10), Inches(5.625)
        )
        bg.fill.solid()
        bg.fill.fore_color.rgb = self._rgb(color_key)
        bg.line.fill.background()
        slide.shapes._spTree.remove(bg._element)
        slide.shapes._spTree.insert(2, bg._element)

    def _add_text(self, slide, x, y, w, h, text, size, bold=False, color='text', align=PP_ALIGN.LEFT, font='body'):
        tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
        tf = tb.text_frame
        tf.text = text
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.font.size = Pt(size)
        p.font.bold = bold
        p.font.name = self.theme[font]
        p.font.color.rgb = self._rgb(color)
        p.alignment = align
        tb.line.fill.background()
        return tb

    # Layout 1: Hero Title
    def layout_01(self):
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self._add_bg(slide, 'primary')

        overlay = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(10), Inches(5.625))
        overlay.fill.solid()
        overlay.fill.fore_color.rgb = self._rgb('secondary')
        overlay.fill.transparency = 0.5
        overlay.line.fill.background()

        self._add_text(slide, 1.5, 1.8, 7, 1.5, "Presentation Title", 60, True, 'bg', PP_ALIGN.CENTER, 'heading')
        self._add_text(slide, 2, 3.5, 6, 0.6, "Subtitle or tagline here", 24, False, 'bg', PP_ALIGN.CENTER, 'body')

    # Layout 2: Section Divider
    def layout_02(self):
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self._add_bg(slide, 'bg')

        self._add_text(slide, 1, 2, 8, 1.2, "Section Title", 48, True, 'primary', PP_ALIGN.CENTER, 'heading')

        bar = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(3.5), Inches(3.5), Inches(3), Inches(0.15))
        bar.fill.solid()
        bar.fill.fore_color.rgb = self._rgb('accent')
        bar.line.fill.background()

    # Layout 3: Content + Visual
    def layout_03(self):
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self._add_bg(slide, 'bg')

        self._add_text(slide, 0.75, 0.75, 8.5, 0.6, "Content Title", 36, True, 'primary', PP_ALIGN.LEFT, 'heading')
        self._add_text(slide, 0.75, 1.75, 4.5, 3, "Main content goes here\n\nKey point one\nKey point two\nKey point three", 18, False, 'text', PP_ALIGN.LEFT, 'body')

        vis = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(5.75), Inches(1.75), Inches(3.5), Inches(3))
        vis.fill.solid()
        vis.fill.fore_color.rgb = self._rgb('bg_alt')
        vis.line.color.rgb = self._rgb('primary')
        vis.line.width = Pt(2)

    # Layout 4: Feature Cards
    def layout_04(self):
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self._add_bg(slide, 'bg_alt')

        self._add_text(slide, 0.75, 0.75, 8.5, 0.6, "Key Features", 36, True, 'primary', PP_ALIGN.CENTER, 'heading')

        for i in range(3):
            x = 1 + i * 2.7
            card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(1.5), Inches(2.5), Inches(2.8))
            card.fill.solid()
            card.fill.fore_color.rgb = self._rgb('bg')
            card.line.fill.background()

            circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x + 0.95), Inches(1.8), Inches(0.6), Inches(0.6))
            circle.fill.solid()
            circle.fill.fore_color.rgb = self._rgb('primary')
            circle.fill.transparency = 0.2
            circle.line.fill.background()

            self._add_text(slide, x + 0.2, 2.6, 2.1, 0.5, f"Feature {i+1}", 20, True, 'text', PP_ALIGN.CENTER, 'heading')
            self._add_text(slide, x + 0.2, 3.2, 2.1, 0.9, "Description of feature", 14, False, 'text_light', PP_ALIGN.CENTER, 'body')

    # Layout 5: Data Chart
    def layout_05(self):
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self._add_bg(slide, 'bg')

        self._add_text(slide, 0.75, 0.75, 8.5, 0.6, "Data Insights", 36, True, 'primary', PP_ALIGN.LEFT, 'heading')

        chart = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.75), Inches(1.75), Inches(8.5), Inches(3.2))
        chart.fill.solid()
        chart.fill.fore_color.rgb = self._rgb('bg_alt')
        chart.line.color.rgb = self._rgb('primary')
        chart.line.width = Pt(1)

    # Layout 6: Quote
    def layout_06(self):
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self._add_bg(slide, 'bg_alt')

        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.5), Inches(1.5), Inches(7), Inches(2.5))
        card.fill.solid()
        card.fill.fore_color.rgb = self._rgb('bg')
        card.line.fill.background()

        self._add_text(slide, 2, 1.9, 6, 1.5, '"Inspiring quote or key message"', 28, False, 'text', PP_ALIGN.CENTER, 'heading')
        self._add_text(slide, 2, 3.5, 6, 0.4, "- Author Name", 16, False, 'text_light', PP_ALIGN.CENTER, 'body')

    # Layout 7: Split Screen
    def layout_07(self):
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])

        left = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(5), Inches(5.625))
        left.fill.solid()
        left.fill.fore_color.rgb = self._rgb('primary')
        left.line.fill.background()

        right = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(5), Inches(0), Inches(5), Inches(5.625))
        right.fill.solid()
        right.fill.fore_color.rgb = self._rgb('bg')
        right.line.fill.background()

        self._add_text(slide, 0.75, 1.5, 3.5, 2.5, "Left Section\n\nContent here", 20, False, 'bg', PP_ALIGN.LEFT, 'body')
        self._add_text(slide, 5.75, 1.5, 3.5, 2.5, "Right Section\n\nContent here", 20, False, 'text', PP_ALIGN.LEFT, 'body')

    # Layout 8: Timeline
    def layout_08(self):
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self._add_bg(slide, 'bg')

        self._add_text(slide, 0.75, 0.75, 8.5, 0.6, "Timeline", 36, True, 'primary', PP_ALIGN.CENTER, 'heading')

        line = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.5), Inches(3), Inches(7), Inches(0.08))
        line.fill.solid()
        line.fill.fore_color.rgb = self._rgb('secondary')
        line.line.fill.background()

        for i in range(4):
            x = 2 + i * 1.8
            circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x - 0.15), Inches(3 - 0.15), Inches(0.3), Inches(0.3))
            circle.fill.solid()
            circle.fill.fore_color.rgb = self._rgb('primary')
            circle.line.fill.background()

            self._add_text(slide, x - 0.6, 2.1, 1.2, 0.4, f"Step {i+1}", 16, True, 'text', PP_ALIGN.CENTER, 'heading')
            self._add_text(slide, x - 0.6, 3.4, 1.2, 0.6, "Description", 12, False, 'text_light', PP_ALIGN.CENTER, 'body')

    # Layout 9: Big Number
    def layout_09(self):
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self._add_bg(slide, 'bg')

        self._add_text(slide, 2, 1.5, 6, 1.8, "95%", 96, True, 'primary', PP_ALIGN.CENTER, 'heading')
        self._add_text(slide, 2, 3.5, 6, 0.8, "Customer Satisfaction Rate", 24, False, 'text', PP_ALIGN.CENTER, 'body')
        self._add_text(slide, 2.5, 4.4, 5, 0.5, "Based on 2024 survey results", 16, False, 'text_light', PP_ALIGN.CENTER, 'body')

    # Layout 10: Image Hero
    def layout_10(self):
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])

        img = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(10), Inches(5.625))
        img.fill.solid()
        img.fill.fore_color.rgb = self._rgb('bg_alt')
        img.line.fill.background()

        overlay = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(10), Inches(5.625))
        overlay.fill.solid()
        overlay.fill.fore_color.rgb = RGBColor(0, 0, 0)
        overlay.fill.transparency = 0.4
        overlay.line.fill.background()

        self._add_text(slide, 1.5, 2, 7, 1.5, "Hero Message", 48, True, 'bg', PP_ALIGN.CENTER, 'heading')

    # Layout 11: Comparison
    def layout_11(self):
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self._add_bg(slide, 'bg')

        self._add_text(slide, 0.75, 0.75, 8.5, 0.6, "Comparison", 36, True, 'primary', PP_ALIGN.CENTER, 'heading')

        for i in range(3):
            x = 1 + i * 2.8
            card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(1.3), Inches(2.6), Inches(3))
            card.fill.solid()
            card.fill.fore_color.rgb = self._rgb('bg_alt')
            card.line.fill.background()

            header = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(1.3), Inches(2.6), Inches(0.5))
            header.fill.solid()
            header.fill.fore_color.rgb = self._rgb('primary')
            header.fill.transparency = 0.15
            header.line.fill.background()

            self._add_text(slide, x + 0.2, 1.45, 2.2, 0.3, f"Option {i+1}", 18, True, 'primary', PP_ALIGN.CENTER, 'heading')
            self._add_text(slide, x + 0.3, 2.1, 2, 1.9, "Feature A\nFeature B\nFeature C", 14, False, 'text', PP_ALIGN.LEFT, 'body')

    # Layout 12: Team Grid
    def layout_12(self):
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self._add_bg(slide, 'bg')

        self._add_text(slide, 0.75, 0.75, 8.5, 0.6, "Our Team", 36, True, 'primary', PP_ALIGN.CENTER, 'heading')

        positions = [(1.2, 1.5), (5.5, 1.5), (1.2, 3.3), (5.5, 3.3)]

        for idx, (x, y) in enumerate(positions):
            card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(3.5), Inches(1.5))
            card.fill.solid()
            card.fill.fore_color.rgb = self._rgb('bg_alt')
            card.line.fill.background()

            circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x + 0.3), Inches(y + 0.35), Inches(0.8), Inches(0.8))
            circle.fill.solid()
            circle.fill.fore_color.rgb = self._rgb('accent')
            circle.fill.transparency = 0.3
            circle.line.fill.background()

            self._add_text(slide, x + 1.3, y + 0.4, 2, 0.4, f"Team Member {idx+1}", 18, True, 'text', PP_ALIGN.LEFT, 'heading')
            self._add_text(slide, x + 1.3, y + 0.9, 2, 0.6, "Job Title", 14, False, 'text_light', PP_ALIGN.LEFT, 'body')

    # Layout 13: Bullet List
    def layout_13(self):
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self._add_bg(slide, 'bg')

        self._add_text(slide, 0.75, 0.75, 8.5, 0.6, "Key Points", 36, True, 'primary', PP_ALIGN.LEFT, 'heading')

        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.5), Inches(1.5), Inches(7), Inches(3))
        card.fill.solid()
        card.fill.fore_color.rgb = self._rgb('bg_alt')
        card.line.fill.background()

        for i in range(4):
            y = 2 + i * 0.6
            bullet = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(2), Inches(y), Inches(0.12), Inches(0.12))
            bullet.fill.solid()
            bullet.fill.fore_color.rgb = self._rgb('accent')
            bullet.line.fill.background()

            self._add_text(slide, 2.3, y - 0.05, 5.9, 0.5, f"Key point {i+1} with detailed information", 18, False, 'text', PP_ALIGN.LEFT, 'body')

    # Layout 14: Call to Action
    def layout_14(self):
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self._add_bg(slide, 'primary')

        overlay = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(10), Inches(5.625))
        overlay.fill.solid()
        overlay.fill.fore_color.rgb = self._rgb('secondary')
        overlay.fill.transparency = 0.6
        overlay.line.fill.background()

        self._add_text(slide, 1.5, 1.5, 7, 1.2, "Ready to Get Started?", 44, True, 'bg', PP_ALIGN.CENTER, 'heading')

        button = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(3.5), Inches(3.2), Inches(3), Inches(0.6))
        button.fill.solid()
        button.fill.fore_color.rgb = self._rgb('accent')
        button.line.fill.background()

        tf = button.text_frame
        tf.text = "Get Started"
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]
        p.font.size = Pt(24)
        p.font.bold = True
        p.font.name = self.theme['heading']
        p.font.color.rgb = self._rgb('bg')
        p.alignment = PP_ALIGN.CENTER

    # Layout 15: Dashboard
    def layout_15(self):
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self._add_bg(slide, 'bg')

        self._add_text(slide, 0.75, 0.75, 8.5, 0.6, "Dashboard Overview", 36, True, 'primary', PP_ALIGN.CENTER, 'heading')

        chart1 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1), Inches(1.5), Inches(4), Inches(3))
        chart1.fill.solid()
        chart1.fill.fore_color.rgb = self._rgb('bg_alt')
        chart1.line.color.rgb = self._rgb('primary')
        chart1.line.width = Pt(1)

        chart2 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(5.3), Inches(1.5), Inches(3.7), Inches(3))
        chart2.fill.solid()
        chart2.fill.fore_color.rgb = self._rgb('bg_alt')
        chart2.line.color.rgb = self._rgb('secondary')
        chart2.line.width = Pt(1)

    def generate_all(self):
        """Generate all 15 layouts"""
        self.layout_01()
        self.layout_02()
        self.layout_03()
        self.layout_04()
        self.layout_05()
        self.layout_06()
        self.layout_07()
        self.layout_08()
        self.layout_09()
        self.layout_10()
        self.layout_11()
        self.layout_12()
        self.layout_13()
        self.layout_14()
        self.layout_15()

        print(f"   > Generated {self.theme['name']} - 15 layouts with content")

    def save(self):
        """Save template"""
        os.makedirs('templates', exist_ok=True)
        filename = f"templates/{self.theme_key}_clean.pptx"
        self.prs.save(filename)
        print(f"   > Saved: {filename}")
        return filename


def main():
    """Generate all 10 clean templates"""
    print("\n" + "="*70)
    print("CLEAN TEMPLATE GENERATOR - Professional Templates with Content")
    print("="*70)
    print("\nGenerating 10 professional templates...\n")

    templates = []
    for i, (theme_key, theme_data) in enumerate(CleanTemplateGenerator.THEMES.items(), 1):
        print(f"[{i}/10] {theme_data['name']}")
        gen = CleanTemplateGenerator(theme_key)
        gen.generate_all()
        path = gen.save()
        templates.append(path)
        print()

    print("="*70)
    print("COMPLETE! Generated 10 professional templates")
    print("="*70)
    print("\nGenerated files:")
    for t in templates:
        print(f"   - {t}")
    print("\nAll templates have:")
    print("   - Proper text content (no weird placeholders)")
    print("   - 15 different professional layouts each")
    print("   - Clean, modern designs")
    print("   - Different color schemes")
    print("\n" + "="*70 + "\n")


if __name__ == "__main__":
    main()
